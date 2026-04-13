"""Product Manager — Smart Error Notebook review material preparation.

Environments: filesystem, email, notion, google_sheets, calendar
2 stages: review material preparation → design director feedback update
9 core checkers + 2 bonus (0 keyword-search)
"""
import re
from datetime import datetime
from pathlib import Path

# ── Constants ─────────────────────────────────────────────────────

BACKLOG_DB_NAME = "lingxi_backlog_q2"

BACKLOG_DB_SCHEMA = {
    "Feature ID": {"title": {}},
    "Title": {"rich_text": {}},
    "Priority": {"select": {"options": [
        {"name": "P0"}, {"name": "P1"}, {"name": "P2"},
    ]}},
    "Status": {"select": {"options": [
        {"name": "pending development"}, {"name": "pending evaluation"},
        {"name": "launched"}, {"name": "needs investigation"},
    ]}},
    "Owner": {"rich_text": {}},
    "Target Version": {"select": {"options": [
        {"name": "v2.4"}, {"name": "v2.5"}, {"name": "v2.6"},
    ]}},
}

INITIAL_BACKLOG_ROWS = [
    {"id": "F-201", "title": "Error Auto Categorization", "priority": "P0",
     "status": "pending development", "owner": "Chen Jie", "version": "v2.5"},
    {"id": "F-202", "title": "Error Redo", "priority": "P0",
     "status": "pending development", "owner": "Chen Jie", "version": "v2.5"},
    {"id": "F-203", "title": "AI Problem-Solving Guidance", "priority": "P2",
     "status": "pending evaluation", "owner": "Chen Jie", "version": "v2.6"},
    {"id": "F-204", "title": "Learning Report", "priority": "P1",
     "status": "launched", "owner": "Chen Jie", "version": "v2.4"},
]

SURVEY_HEADER = ["Feature", "Votes", "Rank"]
SURVEY_ROWS = [
    ["Auto Categorization", "45", "1"],
    ["Error Redo", "38", "2"],
    ["AI Problem-Solving Guidance", "25", "3"],
    ["Learning Report", "18", "4"],
]

CAL_NAME = "lingxi_review"


# ── Helpers ───────────────────────────────────────────────────────

def _notion_title(value: str) -> dict:
    return {"title": [{"text": {"content": value}}]}


def _notion_text(value: str) -> dict:
    return {"rich_text": [{"text": {"content": value}}]}


def _notion_select(value: str) -> dict:
    return {"select": {"name": value}}


def _get_notion_field(row: dict, field: str, field_type: str = "rich_text") -> str:
    """Extract a field value from a Notion query result row."""
    props = row.get("properties", {})
    prop = props.get(field, {})
    if field_type == "title":
        parts = prop.get("title", [])
        return "".join(t.get("plain_text", "") for t in parts)
    elif field_type == "rich_text":
        parts = prop.get("rich_text", [])
        return "".join(t.get("plain_text", "") for t in parts)
    elif field_type == "select":
        sel = prop.get("select", {})
        return sel.get("name", "") if sel else ""
    return ""


async def _find_notion_row(ctx, db_name: str, feature_id: str) -> dict | None:
    """Find a Notion row by Feature ID (title field)."""
    rows = await ctx.notion.query_db(db_name)
    for row in rows:
        fid = _get_notion_field(row, "Feature ID", "title")
        if fid == feature_id:
            return row
    return None


def _parse_xlsx_sheet(ctx, filename: str, sheet_name: str) -> list[dict]:
    """Parse an xlsx sheet from workspace/output/ into list of dicts."""
    path = ctx.workspace / "output" / filename
    if not path.exists():
        return []
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
    except Exception:
        return []
    if sheet_name not in wb.sheetnames:
        return []
    ws = wb[sheet_name]
    rows = list(ws.iter_rows(values_only=True))
    if not rows:
        return []
    headers = [str(h).strip().lower() if h else "" for h in rows[0]]
    result = []
    for row in rows[1:]:
        row_dict = {}
        for j, val in enumerate(row):
            if j < len(headers) and headers[j]:
                row_dict[headers[j]] = str(val).strip() if val is not None else ""
        if any(v for v in row_dict.values()):
            result.append(row_dict)
    return result


def _find_xlsx_row(rows: list[dict], column: str, search: str, exact: bool = False) -> dict | None:
    """Find an xlsx row where column matches search string.

    If exact=True, requires exact match (case-insensitive).
    Otherwise, uses substring match (case-insensitive).
    """
    for row in rows:
        val = row.get(column, "")
        if exact:
            if val.strip().lower() == search.lower():
                return row
        else:
            if search.lower() in val.lower():
                return row
    return None


def _parse_pptx_slide_text(ctx, filename: str, slide_index: int) -> str:
    """Extract all text from a specific slide (0-indexed) in a pptx file."""
    path = ctx.workspace / "output" / filename
    if not path.exists():
        return ""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
    except Exception:
        return ""
    if slide_index >= len(prs.slides):
        return ""
    slide = prs.slides[slide_index]
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                texts.append(para.text)
    return " ".join(texts)


def _normalize(text: str) -> str:
    """Normalize text for comparison: lowercase, collapse whitespace."""
    if not text:
        return ""
    return re.sub(r'[\s\u3000]+', ' ', text.lower().strip())


# ── METADATA ──────────────────────────────────────────────────────

METADATA = {
    "id": "pm_task3",
    "name": "Product Manager Review Material Preparation",
    "category": "project_and_product_manager",
    "environments": ["filesystem", "email", "notion", "google_sheets", "calendar"],
    "timeout_seconds": 600,
    "difficulty": "easy",
    "mm_level": "L3",
    "role": "Xiao Su, Product Manager at Lingxi Education",
    "tags": [
        "product-manager", "review", "excel", "ppt", "multimodal",
        "cross-source-contradiction", "notion", "calendar", "email",
    ],
    "env_config": {
        "email": {
            "users": {
                "xiaosu": {"email": "xiaosu@lingxi.com", "password": "xiaosu_pwd"},
                "zhouming": {"email": "zhouming@lingxi.com", "password": "zhouming_pwd"},
                "lifang": {"email": "lifang@lingxi.com", "password": "lifang_pwd"},
                "chenjie": {"email": "chenjie@lingxi.com", "password": "chenjie_pwd"},
            },
        },
        "google_sheets": {
            "task_id": "pm_task3",
        },
    },
}

PROMPT = "Check your workspace for project materials and prepare the review."


# ── Stage Functions ───────────────────────────────────────────────

async def stage0(ctx):
    """2026-03-19: Review material preparation — consolidate all sources."""
    # 1. Upload all assets (personality .md + input materials)
    await ctx.fs.upload_dir(ctx.task_dir / "assets", "/workspace")

    # 2. Create output directory
    await ctx.fs._sandbox.exec("mkdir -p /workspace/output")

    # 3. Create Notion product backlog database + seed records
    await ctx.notion.create_page("Lingxi Academy Q2 Backlog")
    await ctx.notion.create_database(BACKLOG_DB_NAME, BACKLOG_DB_SCHEMA)
    for rec in INITIAL_BACKLOG_ROWS:
        await ctx.notion.add_database_row(BACKLOG_DB_NAME, {
            "Feature ID": _notion_title(rec["id"]),
            "Title": _notion_text(rec["title"]),
            "Priority": _notion_select(rec["priority"]),
            "Status": _notion_select(rec["status"]),
            "Owner": _notion_text(rec["owner"]),
            "Target Version": _notion_select(rec["version"]),
        })

    # 4. Create Google Sheets survey data
    sheet_info = await ctx.google_sheets.create_spreadsheet("lingxi_survey_2026q1")
    sheet_id = sheet_info["sheet_id"]
    await ctx.google_sheets.update_values(
        sheet_id, "Sheet1!A1:C5",
        [SURVEY_HEADER] + SURVEY_ROWS,
    )

    # 5. Seed historical email (noise — HR team building notice)
    await ctx.email.send_email(
        from_user="zhouming",
        to="xiaosu@lingxi.com",
        subject="March Team Building Event Notice — 3/29 Saturday Afternoon",
        body=(
            "Dear colleagues,\n\n"
            "A spring team building event is scheduled for March 29 (Saturday) "
            "from 14:00 to 17:00. The venue is tentatively set at the park near the office.\n"
            "Please plan your schedule accordingly and try to attend.\n\n"
            "Zhou Ming"
        ),
    )

    # 6. Create calendar for review meetings
    await ctx.calendar.create_calendar(CAL_NAME)

    # 7. Notification
    return {
        "notification": (
            "[2026-03-19 Wednesday] There's a product review meeting tomorrow afternoon. "
            "Help me prepare the materials.\n\n"
            "The workspace has user interview transcript (input/user_interview_teacher.txt), "
            "competitor comparison (input/competitor_comparison.md), and previous meeting minutes "
            "(input/last_review_meeting.md) — please review them all.\n"
            "Survey data is on Google Sheets (lingxi_survey_2026q1), pull that.\n"
            "Also check the product backlog on Notion (lingxi_backlog_q2) for the current status.\n\n"
            "Please do two things:\n"
            "1. Organize the feature spec according to the input/feature_spec_template.xlsx template, "
            "output to output/feature_spec.xlsx\n"
            "2. Create a review PPT based on the input/ppt_template.pptx template, "
            "output to output/product_review.pptx\n\n"
            "If the Notion backlog has any status that needs updating, please handle that too.\n"
            "Schedule the review meeting (at the time decided in the last meeting), "
            "find the attendees from the meeting minutes, and send them a notification.\n\n"
            "Your email is xiaosu@lingxi.com. "
            "CEO: zhouming@lingxi.com. Design Director: lifang@lingxi.com. "
            "Technical Lead: chenjie@lingxi.com.\n"
            "Product backlog is in Notion (database: lingxi_backlog_q2). "
            "Survey data is in Google Sheets (lingxi_survey_2026q1). "
            "Use Google Calendar to schedule the review meeting."
        ),
        "time": "2026-03-19T09:00:00+08:00",
    }


async def stage1(ctx):
    """2026-03-19 afternoon: Design Director Li Fang sends feedback email."""
    # 1. Loud: Li Fang email with feedback
    await ctx.email.send_email(
        from_user="lifang",
        to="xiaosu@lingxi.com",
        subject="Review Material Feedback",
        body=(
            "I looked at the spec table you organized. One thing needs to change:\n"
            "I confirmed with Chen Jie about the AI Problem-Solving Guidance feature — "
            "it's technically feasible for Phase 1 using a RAG approach.\n"
            "Move it from Phase 2 to Phase 1, and change the priority to P1.\n"
            "Also update the timeline slide in the PPT accordingly."
        ),
    )

    # 2. Notification — mentions loud event only
    return {
        "notification": (
            "[2026-03-19 afternoon] Design Director Li Fang sent an email with feedback "
            "that needs changes. Please check your inbox and make the requested adjustments."
        ),
        "time": "2026-03-19T15:00:00+08:00",
    }


# ── Checker Functions ─────────────────────────────────────────────

# -- S0: Review Material Preparation --

async def _s0_competitor_guidance(ctx):
    """Sheet1 'AI Problem-Solving Guidance' competitor_support == '2/3'
    (TiHui✅ + XueBa Notes✅ = 2/3, must read competitor_comparison.md)"""
    rows = _parse_xlsx_sheet(ctx, "feature_spec.xlsx", "Feature List")
    if not rows:
        return False
    row = _find_xlsx_row(rows, "func_id", "F-203", exact=True)
    if not row:
        row = _find_xlsx_row(rows, "func_name", "Guidance")
    if not row:
        return False
    cs = row.get("competitor_support", "").strip()
    return cs == "2/3"


async def _s0_learning_report_issue(ctx):
    """Sheet1 'Learning Report' user_feedback == 'issue'
    (teacher says 'can't find entry point' -> issue, not positive)"""
    rows = _parse_xlsx_sheet(ctx, "feature_spec.xlsx", "Feature List")
    if not rows:
        return False
    row = _find_xlsx_row(rows, "func_id", "F-204", exact=True)
    if not row:
        row = _find_xlsx_row(rows, "func_name", "Learning Report")
    if not row:
        return False
    uf = row.get("user_feedback", "").strip().lower()
    return uf == "issue"


async def _s0_guidance_p2(ctx):
    """Sheet1 'AI Problem-Solving Guidance' priority == 'P2' and version == 'v2.6'
    (user wants it, but meeting minutes say 'tentatively Phase 2' -> follow team decision)"""
    rows = _parse_xlsx_sheet(ctx, "feature_spec.xlsx", "Feature List")
    if not rows:
        return False
    row = _find_xlsx_row(rows, "func_id", "F-203", exact=True)
    if not row:
        row = _find_xlsx_row(rows, "func_name", "Guidance")
    if not row:
        return False
    priority = row.get("priority", "").strip()
    version = row.get("version", "").strip()
    return priority == "P2" and version == "v2.6"


async def _s0_backlog_investigation(ctx):
    """Notion F-204 status changed from 'launched' to 'needs investigation'.
    Cross-modal contradiction: Notion says launched, but teacher can't find entry point."""
    notion_row = await _find_notion_row(ctx, BACKLOG_DB_NAME, "F-204")
    if not notion_row:
        return False
    status = _get_notion_field(notion_row, "Status", "select")
    return status == "needs investigation"


async def _s0_calendar(ctx):
    """Google Calendar has review meeting event: date 2026-03-20, time 14:00, title contains 'review'.
    Agent must extract meeting time from PDF meeting minutes and create calendar event."""
    # Search all calendars the agent may have created
    try:
        cals = await ctx.calendar.list_calendars()
    except Exception:
        return False

    events = []
    for cal in cals:
        try:
            cal_events = await ctx.calendar.get_events(cal["name"])
            events.extend(cal_events)
        except Exception:
            continue

    for event in events:
        summary = event.get("summary", "").lower()
        dtstart = str(event.get("dtstart", ""))
        # Check title contains "review" (case-insensitive)
        has_review = "review" in summary
        # Check date is 2026-03-20
        has_date = "2026-03-20" in dtstart
        # Check time contains 14:00
        has_time = "14:00" in dtstart
        if has_review and has_date and has_time:
            return True
    return False


# -- S1: Design Director Feedback Update --

async def _s1_guidance_upgrade(ctx):
    """Sheet1 'AI Problem-Solving Guidance' priority changed to P1, version to v2.5.
    Event-driven update from Li Fang's email."""
    rows = _parse_xlsx_sheet(ctx, "feature_spec.xlsx", "Feature List")
    if not rows:
        return False
    row = _find_xlsx_row(rows, "func_id", "F-203", exact=True)
    if not row:
        row = _find_xlsx_row(rows, "func_name", "Guidance")
    if not row:
        return False
    priority = row.get("priority", "").strip()
    version = row.get("version", "").strip()
    return priority == "P1" and version == "v2.5"


async def _s1_phase_count(ctx):
    """Sheet2 'Summary' phase1_count == '4', phase2_count == '0'.
    Linkage update: after adjusting F-203 version, summary must also change."""
    rows = _parse_xlsx_sheet(ctx, "feature_spec.xlsx", "Summary")
    if not rows:
        return False
    # Build a map of statistic -> value
    stats = {}
    for row in rows:
        stat_key = row.get("statistic", "").strip().lower()
        stat_val = row.get("value", "").strip()
        if stat_key:
            stats[stat_key] = stat_val
    p1_count = stats.get("phase1_count", "")
    p2_count = stats.get("phase2_count", "")
    return p1_count == "4" and p2_count == "0"


async def _s1_backlog_f203(ctx):
    """Notion F-203: priority=P1, status=pending development, target version=v2.5.
    Multi-tool linkage: not just Excel, must also sync update Notion."""
    notion_row = await _find_notion_row(ctx, BACKLOG_DB_NAME, "F-203")
    if not notion_row:
        return False
    priority = _get_notion_field(notion_row, "Priority", "select")
    status = _get_notion_field(notion_row, "Status", "select")
    version = _get_notion_field(notion_row, "Target Version", "select")
    return priority == "P1" and status == "pending development" and version == "v2.5"


async def _s1_ppt_update(ctx):
    """PPT slide 4 (plan) text references AI Guidance in Phase 1 or v2.5 or P1.
    Multi-artifact linkage: Excel + Notion + PPT all need sync.
    Uses proximity matching to avoid false-positive from unrelated Phase 1 mentions."""
    text = _parse_pptx_slide_text(ctx, "product_review.pptx", 3)  # 0-indexed: slide 4
    normalized = _normalize(text)
    if not normalized:
        return False
    # Check for AI Guidance co-occurring with Phase 1 / P1 / v2.5 within proximity
    guidance_patterns = [
        "guidance", "ai problem", "problem-solving", "problem solving",
    ]
    phase1_patterns = [
        "phase 1", "v2.5", "phase1",
    ]
    # Forward proximity: guidance ... phase1 (within 80 chars)
    for gp in guidance_patterns:
        for pp in phase1_patterns:
            pattern = re.escape(gp) + r".{0,80}" + re.escape(pp)
            if re.search(pattern, normalized):
                return True
            # Reverse proximity: phase1 ... guidance
            pattern_rev = re.escape(pp) + r".{0,80}" + re.escape(gp)
            if re.search(pattern_rev, normalized):
                return True
    # Also accept if "p1" appears near guidance (but "p1" alone is too short for global search)
    for gp in guidance_patterns:
        pattern = re.escape(gp) + r".{0,40}" + r"\bp1\b"
        if re.search(pattern, normalized):
            return True
        pattern_rev = r"\bp1\b" + r".{0,40}" + re.escape(gp)
        if re.search(pattern_rev, normalized):
            return True
    return False


# -- Bonus --

async def _b_ppt_issues_slide(ctx):
    """PPT slide 5 (issues) mentions 'Learning Report' issue."""
    text = _parse_pptx_slide_text(ctx, "product_review.pptx", 4)  # 0-indexed: slide 5
    normalized = _normalize(text)
    if not normalized:
        return False
    return "learning report" in normalized


async def _b_email_notification(ctx):
    """Notification email sent to at least one attendee (zhouming, lifang, or chenjie)."""
    for user_key in ("zhouming", "lifang", "chenjie"):
        try:
            emails = await ctx.email.get_emails(user_key)
        except Exception:
            continue
        for email in emails:
            sender = email.get("from", "")
            if isinstance(sender, dict):
                sender = sender.get("email", "")
            sender = str(sender).lower()
            # Check it's from xiaosu (the agent)
            if "xiaosu" not in sender:
                continue
            # Check subject or body mentions review/meeting
            subject = _normalize(email.get("subject", ""))
            body = _normalize(email.get("body", ""))
            combined = subject + " " + body
            if any(k in combined for k in ["review", "march 20", "3/20", "meeting"]):
                return True
    return False


# ── RUBRIC ────────────────────────────────────────────────────────

RUBRIC = {
    "stage0": [
        {"id": "S0_competitor_guidance", "checker": _s0_competitor_guidance, "weight": 1.5},
        {"id": "S0_learning_report_issue", "checker": _s0_learning_report_issue, "weight": 2.0},
        {"id": "S0_guidance_P2", "checker": _s0_guidance_p2, "weight": 1.5},
        {"id": "S0_backlog_investigation", "checker": _s0_backlog_investigation, "weight": 2.0},
        {"id": "S0_calendar", "checker": _s0_calendar, "weight": 1.5},
    ],
    "stage1": [
        {"id": "S1_guidance_upgrade", "checker": _s1_guidance_upgrade, "weight": 2.0},
        {"id": "S1_phase_count", "checker": _s1_phase_count, "weight": 1.5},
        {"id": "S1_backlog_F203", "checker": _s1_backlog_f203, "weight": 2.0},
        {"id": "S1_ppt_update", "checker": _s1_ppt_update, "weight": 1.5},
    ],
    "final": [
        {"id": "B_PPT_issues_slide", "checker": _b_ppt_issues_slide, "weight": 1.0},
        {"id": "B_email_notification", "checker": _b_email_notification, "weight": 1.0},
    ],
}
